from __future__ import annotations

from docx import Document
from pptx import Presentation
import openpyxl

from ext.document_parser.core.engine_base import BaseEngine
from ext.document_parser.core.parse_result import OutputFormat, ParseResult, PageResult, TableFormat


class DocxEngine(BaseEngine):
    engine_name = "docx"
    supported_formats = [".docx"]

    async def parse(self, file_path: str, options: dict | None = None) -> ParseResult:
        doc = Document(file_path)

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        tables_data = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)

            if table_data:
                headers = table_data[0]
                rows = []
                raw = []

                for i, row in enumerate(table_data):
                    raw.append(row)
                    if i > 0:
                        row_dict = {}
                        for j, header in enumerate(headers):
                            if j < len(row):
                                row_dict[str(header)] = row[j]
                        rows.append(row_dict)

                tables_data.append(
                    TableFormat(
                        headers=headers,
                        rows=rows,
                        raw=raw,
                    ),
                )

        content = "\n\n".join(paragraphs)

        pages_result = [
            PageResult(
                page_number=1,
                content=content,
                tables=tables_data,
                images=[],
                metadata={"paragraph_count": len(paragraphs), "table_count": len(tables_data)},
            ),
        ]

        return ParseResult(
            content=content,
            format=OutputFormat.TEXT,
            pages=pages_result,
            page_count=1,
            structured_data={"tables": [t.model_dump() for t in tables_data]} if tables_data else None,
            metadata={"paragraph_count": len(paragraphs), "table_count": len(tables_data)},
            confidence=0.90,
            engine_used="docx",
        )


class XLSXEngine(BaseEngine):
    engine_name = "xlsx"
    supported_formats = [".xlsx", ".xls"]

    async def parse(self, file_path: str, options: dict | None = None) -> ParseResult:
        wb = openpyxl.load_workbook(file_path, data_only=True)

        sheets_data = []
        all_text = []
        pages_result = []

        for sheet_num, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]

            sheet_data = []
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                sheet_data.append(row_values)

            if sheet_data:
                headers = sheet_data[0]
                rows = []
                raw = []

                for i, row in enumerate(sheet_data):
                    raw.append(row)
                    if i > 0:
                        row_dict = {}
                        for j, header in enumerate(headers):
                            if j < len(row):
                                row_dict[header] = row[j]
                        rows.append(row_dict)

                table_format = TableFormat(
                    headers=headers,
                    rows=rows,
                    raw=raw,
                )

                sheets_data.append(
                    {
                        "sheet_name": sheet_name,
                        "table": table_format,
                    },
                )

                sheet_text = f"## Sheet: {sheet_name}\n\n"
                sheet_text += "\n".join([" | ".join(row) for row in sheet_data])
                all_text.append(sheet_text)

                pages_result.append(
                    PageResult(
                        page_number=sheet_num + 1,
                        content=sheet_text,
                        tables=[table_format],
                        images=[],
                        metadata={"sheet_name": sheet_name},
                    ),
                )

        return ParseResult(
            content="\n\n".join(all_text),
            format=OutputFormat.TEXT,
            pages=pages_result,
            page_count=len(sheets_data),
            structured_data={"sheets": sheets_data} if sheets_data else None,
            metadata={"sheet_count": len(sheets_data)},
            confidence=0.95,
            engine_used="xlsx",
        )


class PPTXEngine(BaseEngine):
    engine_name = "pptx"
    supported_formats = [".pptx"]

    async def parse(self, file_path: str, options: dict | None = None) -> ParseResult:
        prs = Presentation(file_path)

        pages_result = []
        all_text = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            page_content = "\n\n".join(slide_text)
            all_text.append(page_content)

            pages_result.append(
                PageResult(
                    page_number=slide_num + 1,
                    content=page_content,
                    tables=[],
                    images=[],
                    metadata={"slide_number": slide_num + 1},
                ),
            )

        return ParseResult(
            content="\n\n".join(all_text),
            format=OutputFormat.TEXT,
            pages=pages_result,
            page_count=len(prs.slides),
            metadata={"slide_count": len(prs.slides)},
            confidence=0.85,
            engine_used="pptx",
        )
